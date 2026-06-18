#!/usr/bin/env python3
"""
SCAIT Presentation Generator
Run: python3 build_scait.py
Output: SCAIT.pptx (same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE

# ── Colors ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x13, 0x19, 0x21)
CARD   = RGBColor(0x1E, 0x2A, 0x35)
CARD2  = RGBColor(0x0D, 0x13, 0x1B)   # darker panel (split slide right side)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0xD1, 0xD5, 0xDB)
MUTED  = RGBColor(0x9C, 0xA3, 0xAF)
BLUE   = RGBColor(0x2D, 0x74, 0xDA)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
BORDER = RGBColor(0x26, 0x33, 0x42)   # ≈ rgba(255,255,255,0.07) on CARD

# ── Dimensions (16:9 — 33.87 × 19.05 cm) ──────────────────────────────────────
W, H  = Inches(13.33), Inches(7.5)
ML    = Inches(0.83)    # left margin
MT    = Inches(0.63)    # top margin
CW    = W - ML * 2      # content width
FONT  = "Inter"


# ── Primitive helpers ──────────────────────────────────────────────────────────

def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = NAVY
    return s


def txt(s, text, x, y, w, h, sz=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text            = text
    r.font.name       = FONT
    r.font.size       = Pt(sz)
    r.font.bold       = bold
    r.font.color.rgb  = color
    return box


def multiline_txt(s, lines, x, y, w, h, sz=52, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER):
    """Multi-paragraph text box (for hero headings with 3 lines)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text           = line
        r.font.name      = FONT
        r.font.size      = Pt(sz)
        r.font.bold      = bold
        r.font.color.rgb = color
    return box


def rnd(s, x, y, w, h, fill=CARD, line_color=BORDER, radius=0.05):
    shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb      = line_color
    shape.line.width          = Pt(0.75)
    shape.adjustments[0]      = radius
    return shape


def dot(s, x, y, d=Inches(0.37), fill=BLUE):
    shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def flat_rect(s, x, y, w, h, fill=NAVY):
    shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def line(s, x1, y1, x2, y2, color=WHITE, width=1.5):
    conn = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width     = Pt(width)
    return conn


# ── Card component ─────────────────────────────────────────────────────────────

def draw_card(s, x, y, w, h, title, body):
    """Product card: rounded rect + blue dot + title + body."""
    rnd(s, x, y, w, h)

    pad   = Inches(0.28)
    dot_d = Inches(0.36)

    dot(s, x + pad, y + pad, dot_d)

    ty = y + pad + dot_d + Inches(0.14)
    txt(s, title,
        x + pad, ty, w - pad * 2, Inches(0.4),
        sz=17, bold=True)

    by = ty + Inches(0.44)
    txt(s, body,
        x + pad, by, w - pad * 2, h - (by - y) - pad,
        sz=12, color=MUTED)


def layer_header(s, title, tagline):
    txt(s, title,   ML, MT,               CW, Inches(0.7),  sz=38, bold=True)
    txt(s, tagline, ML, MT + Inches(0.78), CW, Inches(0.35), sz=13, color=SOFT)


def two_cards(s, t1, b1, t2, b2):
    cy = MT + Inches(1.32)
    ch = H - cy - Inches(0.45)
    cw = (CW - Inches(0.25)) / 2
    draw_card(s, ML,                     cy, cw, ch, t1, b1)
    draw_card(s, ML + cw + Inches(0.25), cy, cw, ch, t2, b2)


def one_card(s, t, b):
    cy = MT + Inches(1.32)
    ch = H - cy - Inches(0.45)
    draw_card(s, ML, cy, CW, ch, t, b)


# ── Slide 1 — Cover ────────────────────────────────────────────────────────────

def s01_cover(prs):
    s = new_slide(prs)
    txt(s, "SCAIT",
        ML, H / 2 - Inches(0.9), CW, Inches(0.9),
        sz=96, bold=True, align=PP_ALIGN.CENTER)
    txt(s, "Supply Chain Artificial Intelligence Technologies",
        ML, H / 2 + Inches(0.1), CW, Inches(0.4),
        sz=16, color=SOFT, align=PP_ALIGN.CENTER)


# ── Slide 2 — Statement: Manual Trap ──────────────────────────────────────────

def s02_manual_trap(prs):
    s = new_slide(prs)
    txt(s, "If your team is drowning in the 'Manual Trap'.",
        ML, H / 2 - Inches(1.0), CW, Inches(0.9),
        sz=42, bold=True, align=PP_ALIGN.CENTER)

    # "70%" is bold, rest is regular — two runs in one paragraph
    box = s.shapes.add_textbox(ML, H / 2 + Inches(0.05), CW, Inches(0.55))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    r1 = p.add_run()
    r1.text = "70%"
    r1.font.name = FONT; r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = WHITE

    r2 = p.add_run()
    r2.text = (" of Senior Management's time is wasted on manual Excel updates"
               " and messaging apps. This isn't leadership; it's firefighting.")
    r2.font.name = FONT; r2.font.size = Pt(16); r2.font.bold = False
    r2.font.color.rgb = SOFT


# ── Slide 3 — Split: Invisible Gaps ───────────────────────────────────────────

def s03_gaps(prs):
    s = new_slide(prs)

    # Right dark panel (texture placeholder)
    flat_rect(s, W / 2, Inches(0), W / 2, H, fill=CARD2)

    txt(s, "You are losing money in invisible gaps.",
        ML, H / 2 - Inches(1.1), W / 2 - ML - Inches(0.2), Inches(1.3),
        sz=34, bold=True)
    txt(s,
        ("3PL overcharges, supplier price creep, and quality deviations"
         " erode up to 5% of your net margin. If you aren't auditing 100%"
         " of your data, you are overpaying."),
        ML, H / 2 + Inches(0.35), W / 2 - ML - Inches(0.2), Inches(1.1),
        sz=14, color=SOFT)


# ── Slide 4 — Statement: One Architecture ─────────────────────────────────────

def s04_architecture(prs):
    s = new_slide(prs)
    multiline_txt(s,
        ["SCAIT", "One Architecture.", "Zero Friction."],
        ML, H / 2 - Inches(1.55), CW, Inches(2.85),
        sz=52, bold=True, align=PP_ALIGN.CENTER)
    txt(s,
        ("We don't just provide software. We deploy an autonomous intelligence"
         " layer that sees, analyzes, and acts on your behalf."),
        ML, H / 2 + Inches(1.1), CW, Inches(0.5),
        sz=14, color=SOFT, align=PP_ALIGN.CENTER)


# ── Slide 5 — Layer 1: Planning & Analytics ────────────────────────────────────

def s05_layer1(prs):
    s = new_slide(prs)
    layer_header(s, "Layer 1: Planning & Analytics",
                 "Strategic Core: Real-time control over inventory and risks.")
    two_cards(s,
        "Demand vs Plan Deviation Monitor",
        ("Automated tracking of gaps between Sales Forecasts and Supply Plans."
         " Instant alerts when deviations exceed 5% of target KPIs."
         " Prevention of Stock-outs and capital tied up in Overstock."),
        "Exception Management System",
        ("A single 'Control Tower' for all critical events (delays, defects,"
         " price changes). AI-driven prioritization of issues based on their"
         " financial impact on P&L."))


# ── Slide 6 — Layer 2: Procurement & Quality ──────────────────────────────────

def s06_layer2a(prs):
    s = new_slide(prs)
    layer_header(s, "Layer 2: Procurement & Quality",
                 "Operational Core: Eliminating the manual follow-up loop.")
    two_cards(s,
        "Supplier Communication AI",
        ("24/7 autonomous correspondence with factories regarding statuses"
         " and docs. Saves 70% of Procurement's time by automating the"
         " routine follow-up."),
        "AI Product Spec & Inspection Generator",
        ("Instant generation of technical Tech Packs and Inspection Sheets."
         " Standardization of quality requirements to eliminate production errors."))


# ── Slide 7 — Layer 2: Quality & Validation ───────────────────────────────────

def s07_layer2b(prs):
    s = new_slide(prs)
    layer_header(s, "Layer 2: Quality & Validation",
                 "Security Core: Guarding the 'Golden Standard.'")
    two_cards(s,
        "AI Quality Inspection Validation System",
        ("Automated verification of inspector reports to detect fraud or errors."
         " Instant Pass/Fail verdicts based on AI analysis of photos and data."),
        "AI Inspection Report Processor",
        ("Digitizes unstructured PDF/Excel and QC photos reports into a searchable"
         " database. Predictive analytics to identify defect trends per supplier."))


# ── Slide 8 — Layer 3: Logistics & Tracking ───────────────────────────────────

def s08_layer3(prs):
    s = new_slide(prs)
    layer_header(s, "Layer 3: Logistics & Tracking",
                 "Movement Core: Total visibility of the floating capital.")
    one_card(s,
        "Forwarder Communication & Shipment Tracker AI",
        ("Real-time container tracking via carrier APIs (No manual checks)."
         " Automated ETA/ETD updates and document synchronization with forwarders."
         " AI-driven rescheduling scenarios for port delays or route changes."))


# ── Slide 9 — Layer 4: Financial Audit ────────────────────────────────────────

def s09_layer4(prs):
    s = new_slide(prs)
    layer_header(s, "Layer 4: Financial Audit",
                 "Profit Core: The ROI Center.")
    two_cards(s,
        "Supplier Invoice Validation System",
        ("Instant cross-check of Invoices vs. POs (price, quantity, terms)."
         " Detection of hidden surcharges or unauthorized price hikes."),
        "3PL Invoice & Cost Validation System",
        ("Deep audit of warehouse and last-mile billing (DIM weight, picking,"
         " storage). Direct EBITDA impact through automated recovery of overcharges."))


# ── Slide 10 — SCAIT Core Diagram ─────────────────────────────────────────────

def s10_diagram(prs):
    s = new_slide(prs)

    txt(s, "A Unified Neural Network for your Supply Chain.",
        ML, MT, CW, Inches(0.6),
        sz=28, bold=True, align=PP_ALIGN.CENTER)
    txt(s,
        ("Every Layer feeds the next. Financial data validates Logistics,"
         " Quality data refines Procurement."
         " Seamless API integration with your existing ERP/WMS."),
        ML, MT + Inches(0.65), CW, Inches(0.45),
        sz=12, color=SOFT, align=PP_ALIGN.CENTER)

    # Diagram bounds
    da_x = ML + Inches(0.6)
    da_y = MT + Inches(1.3)
    da_w = CW - Inches(1.2)
    da_h = H - da_y - Inches(0.4)

    # Center oval — SCAIT Core
    cw, ch = Inches(1.9), Inches(1.15)
    cx = da_x + (da_w - cw) / 2
    cy = da_y + (da_h - ch) / 2
    ccx, ccy = cx + cw / 2, cy + ch / 2

    core = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, cx, cy, cw, ch)
    core.fill.solid()
    core.fill.fore_color.rgb = WHITE
    core.line.color.rgb = SOFT
    core.line.width = Pt(1.5)

    # Two-line label inside oval
    oval_box = s.shapes.add_textbox(cx, cy + Inches(0.1), cw, ch - Inches(0.2))
    oval_tf  = oval_box.text_frame
    oval_tf.word_wrap = False
    for i, line_text in enumerate(["SCAIT", "Core"]):
        p = oval_tf.paragraphs[0] if i == 0 else oval_tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line_text
        r.font.name = FONT; r.font.size = Pt(14)
        r.font.bold = True; r.font.color.rgb = NAVY

    # Corner boxes — added AFTER connectors so they render on top
    bw, bh = Inches(2.4), Inches(1.5)
    boxes = [
        ("Financial\nAudit",       da_x,                da_y),
        ("Planning &\nAnalytics",  da_x + da_w - bw,    da_y),
        ("Logistics &\nTracking",  da_x,                da_y + da_h - bh),
        ("Procurement &\nQuality", da_x + da_w - bw,    da_y + da_h - bh),
    ]

    # Draw connectors first (they go behind)
    for _, bx, by in boxes:
        bcx, bcy = bx + bw / 2, by + bh / 2
        line(s, bcx, bcy, ccx, ccy, color=SOFT, width=1.5)

    # Draw boxes on top
    for label, bx, by in boxes:
        rnd(s, bx, by, bw, bh, fill=CARD, line_color=WHITE, radius=0.06)
        lbox = s.shapes.add_textbox(
            bx + Inches(0.15), by + Inches(0.3),
            bw - Inches(0.3), bh - Inches(0.45))
        ltf = lbox.text_frame
        ltf.word_wrap = True
        for i, lline in enumerate(label.split("\n")):
            p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = lline
            r.font.name = FONT; r.font.size = Pt(13)
            r.font.bold = True; r.font.color.rgb = WHITE

    # Redraw oval + label on top of everything
    top_oval = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, cx, cy, cw, ch)
    top_oval.fill.solid()
    top_oval.fill.fore_color.rgb = WHITE
    top_oval.line.color.rgb = SOFT
    top_oval.line.width = Pt(1.5)

    top_label = s.shapes.add_textbox(cx, cy + Inches(0.1), cw, ch - Inches(0.2))
    top_ltf = top_label.text_frame
    top_ltf.word_wrap = False
    for i, line_text in enumerate(["SCAIT", "Core"]):
        p = top_ltf.paragraphs[0] if i == 0 else top_ltf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line_text
        r.font.name = FONT; r.font.size = Pt(14)
        r.font.bold = True; r.font.color.rgb = NAVY


# ── Slide 11 — Stats ───────────────────────────────────────────────────────────

def s11_stats(prs):
    s = new_slide(prs)

    txt(s, "Efficiency is no longer an option.",
        ML, MT + Inches(0.25), CW, Inches(0.65),
        sz=36, bold=True, align=PP_ALIGN.CENTER)

    col_w = CW / 3
    stats = [
        ("80%",  "reduction in manual\ncommunication."),
        ("100%", "audit coverage of all\nfinancial transactions."),
        ("0",    "missed deviations in\nthe supply chain."),
    ]
    ny = H * 0.38

    for i, (num, label) in enumerate(stats):
        x = ML + col_w * i
        txt(s, num,
            x, ny, col_w, Inches(1.1),
            sz=80, bold=True, align=PP_ALIGN.CENTER)
        # Two-line label
        label_box = s.shapes.add_textbox(x, ny + Inches(1.15), col_w, Inches(0.7))
        label_tf  = label_box.text_frame
        label_tf.word_wrap = True
        for j, lline in enumerate(label.split("\n")):
            p = label_tf.paragraphs[0] if j == 0 else label_tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = lline
            r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = SOFT


# ── Slide 12 — CTA ─────────────────────────────────────────────────────────────

def s12_cta(prs):
    s = new_slide(prs)

    txt(s, "Let's build the future of your supply chain today",
        ML, Inches(1.2), CW, Inches(0.75),
        sz=32, bold=True, align=PP_ALIGN.CENTER)
    txt(s,
        ("Click the link above to schedule a 20-minute session.\n"
         "We will define your automation roadmap and identify"
         " your most immediate ROI opportunities."),
        ML, Inches(2.1), CW, Inches(0.75),
        sz=14, color=SOFT, align=PP_ALIGN.CENTER)

    contacts = [
        ("Artem Stepanenko | CEO/Founder", SOFT),
        ("Email: artem@scait.space",       ORANGE),
        ("Phone: +380686868880",           SOFT),
        ("www.scait.space",                ORANGE),
        ("https://calendly.com/artem-scait/30min/", SOFT),
    ]
    cy = Inches(3.3)
    for text, color in contacts:
        txt(s, text, ML, cy, CW, Inches(0.35),
            sz=14, color=color, align=PP_ALIGN.CENTER)
        cy += Inches(0.38)


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s01_cover(prs)
    s02_manual_trap(prs)
    s03_gaps(prs)
    s04_architecture(prs)
    s05_layer1(prs)
    s06_layer2a(prs)
    s07_layer2b(prs)
    s08_layer3(prs)
    s09_layer4(prs)
    s10_diagram(prs)
    s11_stats(prs)
    s12_cta(prs)

    import os
    here = os.path.dirname(os.path.abspath(__file__))
    out  = os.path.join(here, "SCAIT.pptx")
    prs.save(out)
    print(f"✓  Saved  → {out}")
    print(f"   Slides : {len(prs.slides)}")


if __name__ == "__main__":
    build()
